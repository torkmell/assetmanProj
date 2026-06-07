"""Build the 15-slide pitch deck (PowerPoint) for GSD2T Asset Management, flagship = V1 defensive sleeve.
Reads v1_flagship.json (+ sectorwide_full.json, survivorship_corrected.json, bakeoff_*.json).
Output: GSD2T_Pitch_Deck.pptx  (export to PDF from PowerPoint for the <=15-slide submission)
"""
import json
from pathlib import Path
import numpy as np, pandas as pd
import matplotlib; matplotlib.use("Agg"); import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

V=json.load(open("v1_flagship.json")); F=json.load(open("sectorwide_full.json"))
PIT=json.load(open("survivorship_corrected.json")); BV=json.load(open("bakeoff_variants.json")); BW=json.load(open("bakeoff_improvements.json"))
CV=json.load(open("concentration_v1.json")); CS=json.load(open("concentration_survivorship.json"))
RB=json.load(open("robustness_assumptions.json"))
S=V["summary"]["Fund (full)"]; ISm=V["summary"]["Fund IS (2002-2015)"]; OOS=V["summary"]["Fund OOS (2016-2026)"]
SPY=V["summary"]["SPY"]; AL=V["alpha"]; NF=V["net_of_fee"]; CAP=F["capacity"]; SEN=F["sensitivity"]
SEC=F["sector_breakdown"]; HL=F["holdings_latest"]; BZ=PIT["bias"]

NAVY=RGBColor(0x0B,0x1F,0x3A); NAVYL=RGBColor(0x1E,0x3A,0x6C); GOLD=RGBColor(0xC9,0xA9,0x6E)
CREAM=RGBColor(0xF5,0xF2,0xEA); GREY=RGBColor(0x5A,0x6F,0x8C); POS=RGBColor(0x2E,0x7D,0x5B)
NEG=RGBColor(0xB8,0x3A,0x3A); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x22,0x2A,0x38)

# ---------- deck charts ----------
def ser(p): i=pd.to_datetime([x[0] for x in p]); return pd.Series([x[1] for x in p],index=i)
fund=ser(V["equity_curves"]["Fund (V1)"]); spy=ser(V["equity_curves"]["SPY"])
plt.rcParams.update({"font.family":"DejaVu Sans","font.size":13})
fig,ax=plt.subplots(figsize=(7.6,4.3))
ax.plot(fund.index,fund.values,color="#0B1F3A",lw=2.6,label=f"GSD2T (V1)  {S['CAGR']*100:.1f}% CAGR · Sharpe {S['Sharpe']:.2f}")
ax.plot(spy.index,spy.values,color="#9aa6b5",lw=1.6,ls="--",label=f"S&P 500  {SPY['CAGR']*100:.1f}% · Sharpe {SPY['Sharpe']:.2f}")
ax.set_yscale("log"); ax.set_ylabel("Growth of $1 (log)"); ax.legend(fontsize=10,loc="upper left"); ax.grid(alpha=.25)
ax.set_title("Simulated growth of $1 — net of trading costs (2002–2026)",fontsize=12,color="#0B1F3A")
fig.tight_layout(); fig.savefig("fig_deck_equity.png",dpi=150); plt.close(fig)

st=V["stress"]; fig,ax=plt.subplots(figsize=(7.6,4.0)); x=np.arange(len(st)); ww=0.4
ax.bar(x-ww/2,[s["fund"]*100 for s in st],ww,label="GSD2T",color="#0B1F3A")
ax.bar(x+ww/2,[s["spy"]*100 for s in st],ww,label="S&P 500",color="#C9A96E")
ax.set_xticks(x); ax.set_xticklabels([s["window"] for s in st],rotation=25,ha="right",fontsize=9)
ax.axhline(0,color="#444",lw=.8); ax.set_ylabel("Total return %"); ax.legend(fontsize=10)
ax.set_title("Crisis stress tests — total return through each window",fontsize=12,color="#0B1F3A")
fig.tight_layout(); fig.savefig("fig_deck_stress.png",dpi=150); plt.close(fig)

# ---------- pptx helpers ----------
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); BLANK=prs.slide_layouts[6]
SW,SH=13.333,7.5; _n=[0]
def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def rect(s,x,y,w,h,color,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid()
    sh.fill.fore_color.rgb=color; sh.shadow.inherit=False
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    return sh
def txt(s,x,y,w,h,runs,size=14,color=NAVY,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4,font="Calibri"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[runs]
    for i,item in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(space)
        segs=item if isinstance(item,list) else [(item,bold,color)]
        for seg in segs:
            t,b,c=(seg if isinstance(seg,tuple) and len(seg)==3 else (seg,bold,color))
            r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=b; r.font.color.rgb=c; r.font.name=font
    return tb
def bullets(s,x,y,w,h,items,size=15,color=DARK,gap=7,mark="▪  ",markcolor=GOLD):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        rm=p.add_run(); rm.text=mark; rm.font.size=Pt(size); rm.font.color.rgb=markcolor; rm.font.bold=True; rm.font.name="Calibri"
        if isinstance(it,tuple):
            lead,rest=it; r1=p.add_run(); r1.text=lead; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=NAVY; r1.font.name="Calibri"
            r2=p.add_run(); r2.text=rest; r2.font.size=Pt(size); r2.font.color.rgb=color; r2.font.name="Calibri"
        else:
            r=p.add_run(); r.text=it; r.font.size=Pt(size); r.font.color.rgb=color; r.font.name="Calibri"
    return tb
def header(s,title,kicker=None):
    rect(s,0,0,SW,1.18,NAVY); rect(s,0,1.18,SW,0.06,GOLD)
    txt(s,0.55,0.30,11.5,0.7,title,size=26,color=WHITE,bold=True)
    if kicker: txt(s,0.57,0.86,11.5,0.3,kicker,size=12,color=GOLD,bold=True)
    _n[0]+=1; txt(s,12.4,0.45,0.7,0.4,f"{_n[0]:02d}",size=12,color=GOLD,bold=True,align=PP_ALIGN.RIGHT)
    txt(s,0.55,7.08,9,0.3,"GSD2T ASSET MANAGEMENT · SIMULATED / FICTIONAL PITCH",size=8,color=GREY)
def table(s,x,y,w,rowh,headers,rows,colw=None,fs=10.5,hfill=NAVY):
    nr,nc=len(rows)+1,len(headers); from pptx.util import Inches as I
    gt=s.shapes.add_table(nr,nc,I(x),I(y),I(w),I(rowh*nr)).table
    if colw:
        tot=sum(colw)
        for j,cw in enumerate(colw): gt.columns[j].width=I(w*cw/tot)
    for j,hh in enumerate(headers):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=hfill; c.margin_top=Pt(2); c.margin_bottom=Pt(2)
        p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=hh; r.font.size=Pt(fs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name="Calibri"
        p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=gt.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=CREAM if i%2 else WHITE; c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=str(val); r.font.size=Pt(fs); r.font.name="Calibri"
            r.font.color.rgb=NAVY; r.font.bold=(j==0)
            p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
    return gt
def kpis(s,y,cards,h=1.25,w=2.75,gap=0.32,x0=0.55):
    for i,(lab,val,sub) in enumerate(cards):
        x=x0+i*(w+gap); rect(s,x,y,w,h,NAVY); rect(s,x,y,0.08,h,GOLD)
        txt(s,x+0.2,y+0.12,w-0.3,0.3,lab,size=10,color=GOLD,bold=True)
        txt(s,x+0.2,y+0.36,w-0.3,0.55,val,size=26,color=WHITE,bold=True,font="Georgia")
        txt(s,x+0.2,y+0.92,w-0.3,0.3,sub,size=10,color=RGBColor(0xB8,0xC8,0xE0))
def pc(x): return f"{x*100:.1f}%"

# ===== 1. TITLE =====
s=slide(NAVY); rect(s,0,3.18,SW,0.05,GOLD)
txt(s,0.9,2.0,11.5,1.0,[[("GSD",True,WHITE),("2",True,GOLD),("T  ",True,WHITE),("Asset Management",True,WHITE)]],size=46,font="Georgia")
txt(s,0.92,3.35,11.5,0.5,"A diversified, macro-aware equity strategy that compounds by losing less.",size=18,color=GOLD)
txt(s,0.92,4.15,11.5,0.4,"Quant Fund Pitch · ESADE Asset Management · seeking a USD 100M commitment",size=14,color=WHITE)
txt(s,0.92,6.5,11.5,0.4,"SIMULATED / FICTIONAL pitch for academic evaluation. Past (simulated) performance does not indicate future results.",size=10,color=GREY)

# ===== 2. EXECUTIVE SUMMARY =====
s=slide(); header(s,"Executive summary","THE ASK · USD 100M")
kpis(s,1.55,[("NET CAGR (TO INVESTOR)",pc(NF["net"]["CAGR"]),f"vs S&P 500 {pc(NF['spy']['CAGR'])}"),
             ("SHARPE (GROSS)",f"{S['Sharpe']:.2f}",f"vs S&P 500 {SPY['Sharpe']:.2f}"),
             ("MAX DRAWDOWN",pc(S['MaxDD']),f"vs S&P 500 {pc(SPY['MaxDD'])}"),
             ("ALPHA (FF5+MOM)",f"+{AL['annualized']*100:.1f}%",f"t = {AL['tstat']:.1f}")])
bullets(s,0.55,3.25,12.2,3.4,[
 ("What we are: ","a long-only, systematic S&P 500 strategy with a macro-regime risk overlay and a productive defensive sleeve. Daily-liquid, transparent, factor-aware."),
 ("The edge: ","we do not claim a faster momentum signal — we claim superior capital preservation. A 4-factor macro overlay cuts equity exposure in stress and rotates into Treasuries + gold, not idle cash."),
 ("The result (simulated, 2002–2026): ","market-beating compounding at ~2/3 of the market's volatility and less than half its drawdown; statistically significant alpha that holds out-of-sample (IS Sharpe %.2f → OOS %.2f)."%(ISm['Sharpe'],OOS['Sharpe'])),
 ("Honesty first: ","survivorship bias measured on point-in-time Bloomberg data (~1%/yr), full transaction costs, train/test discipline, and standard enhancements tested and rejected."),
],size=14,gap=9)

# ===== 3. THE EDGE =====
s=slide(); header(s,"The edge — what inefficiency, and why it persists","EDGE & STRATEGY COHERENCE")
bullets(s,0.55,1.55,6.0,5.0,[
 ("The inefficiency: ","investors systematically overpay in calm markets and panic-sell in stress. Cross-sectional momentum persists because of under-reaction; drawdowns persist because risk premia are time-varying."),
 ("Our claim: ","the durable money is not in a better stock signal (commoditised) — it is in disciplined, rules-based de-risking that most discretionary investors cannot execute under pressure."),
 ("Why it persists: ","it is behavioural and structural, not arbitrageable away — mandates, career risk and leverage constraints stop most managers from cutting exposure when signals say so."),
 ("Why us: ","a transparent, repeatable process with no discretion, no leverage, and a defensive sleeve that turns risk-off periods from a drag into a contributor."),
],size=14,gap=11)
rect(s,7.0,1.7,5.7,4.6,CREAM)
txt(s,7.25,1.9,5.2,0.4,"WHERE THE EDGE SHOWS UP",size=12,color=NAVY,bold=True)
table(s,7.25,2.4,5.2,0.62,["Metric","GSD2T","S&P 500"],
 [["Sharpe (risk-adjusted)",f"{S['Sharpe']:.2f}",f"{SPY['Sharpe']:.2f}"],
  ["Max drawdown",pc(S['MaxDD']),pc(SPY['MaxDD'])],
  ["Volatility",pc(S['Vol']),pc(SPY['Vol'])],
  ["Worst crisis (GFC)","see stress","−50%+"]],colw=[2.2,1,1],fs=11)
txt(s,7.25,5.5,5.2,0.7,"We win on the denominator (risk), not by chasing the numerator (return).",size=12,color=NAVY,bold=True)

# ===== 4. DATA & SIGNALS =====
s=slide(); header(s,"Data & signals","DATA HYGIENE · NO LOOK-AHEAD")
txt(s,0.55,1.4,12,0.4,"Two signal layers, both computed only from information available at decision time (lagged), then traded the next month.",size=13,color=GREY)
txt(s,0.55,2.0,5.9,0.35,"LAYER 1 — STOCK SELECTION",size=12,color=NAVY,bold=True)
table(s,0.55,2.4,5.9,0.52,["Item","Detail"],
 [["Signal","12-1 month price momentum"],["Universe","Full S&P 500, all 11 sectors"],
  ["Ranking","Cross-sectional z-score each month"],["Selection","Equal-weight top quartile (~125 names)"]],colw=[1,2],fs=10.5)
txt(s,6.8,2.0,6.0,0.35,"LAYER 2 — MACRO OVERLAY (4 FACTORS)",size=12,color=NAVY,bold=True)
table(s,6.8,2.4,6.0,0.52,["Factor","Captures","Source"],
 [["VIX","Equity fear","CBOE"],["Credit (IEF−HYG)","Credit stress","ETF px"],
  ["10Y yield Δ","Rate tightening","UST"],["Index trend","Market momentum","S&P 500"]],colw=[1.4,1.3,0.9],fs=10)
rect(s,0.55,5.05,12.2,1.5,CREAM)
bullets(s,0.75,5.15,12.0,1.4,[
 ("Data sources: ","yfinance (prices, VIX, ETFs), Ken French library (FF5 + Momentum, risk-free), Bloomberg point-in-time S&P 500 membership for the survivorship test."),
 ("Look-ahead controls: ","all signals lagged ≥1 month; rolling z-scores use only trailing windows; 2000–2001 reserved as signal warm-up, performance measured from 2002."),
],size=12,gap=7)

# ===== 5. STRATEGY =====
s=slide(); header(s,"The strategy — two layers plus a defensive sleeve","STRATEGY")
for i,(t,d,col) in enumerate([
 ("1 · SELECT","Buy the strongest momentum stocks across all 11 S&P 500 sectors. Equal-weight, diversified — not a tech bet.",NAVY),
 ("2 · OVERLAY","A 4-factor macro 'risk dial' scales equity exposure 30–100%. Calm → fully invested; stress → de-risk hard.",NAVYL),
 ("3 · DEFENSIVE SLEEVE","The de-risked portion earns Treasuries + gold (convexity in crises), not idle cash. This is V1 — our key refinement.",POS)]):
    x=0.55+i*4.15; rect(s,x,1.6,3.9,2.6,col); txt(s,x+0.25,1.8,3.5,0.4,t,size=15,color=GOLD,bold=True)
    txt(s,x+0.25,2.45,3.5,1.6,d,size=12.5,color=WHITE)
txt(s,0.55,4.5,12.2,0.4,"Position per stock = (1 / number held) × macro exposure.   No leverage · no shorting · monthly rebalance · 15 bps round-trip cost.",size=13,color=NAVY,bold=True)
bullets(s,0.55,5.1,12.2,1.4,[
 ("Genuinely diversified: ","currently %d holdings across %d sectors — concentration risk is structurally capped (~1%% per name)."%(HL['n_holdings'],HL['n_sectors'])),
 ("It generalises: ","the same engine gives ~1.0 Sharpe on the Dow 30 and the tech sub-universe — evidence of a real effect, not a single-market fit."),
],size=13,gap=7)

# ===== 6. PORTFOLIO CONSTRUCTION =====
s=slide(); header(s,"Portfolio construction","PORTFOLIO CONSTRUCTION")
txt(s,0.55,1.45,6.0,0.35,"CURRENT SECTOR SPREAD (DIVERSIFIED)",size=12,color=NAVY,bold=True)
secrows=[[k,str(v)] for k,v in list(SEC.items())[:8]]
table(s,0.55,1.85,5.6,0.40,["Sector","# stocks"],secrows,colw=[2.4,1],fs=10)
txt(s,6.8,1.45,6.0,0.35,"CONSTRUCTION RULES",size=12,color=NAVY,bold=True)
bullets(s,6.8,1.9,6.0,4.4,[
 ("Weighting: ","equal-weight within the book (tested inverse-vol and concentration — neither improved risk-adjusted returns)."),
 ("Exposure: ","macro overlay sets gross 30–100%%; average ~%d%% in equities, the rest in the defensive sleeve."%((1-V['ops']['avg_cash'])*100)),
 ("Defensive sleeve: ","Treasuries + gold (TLT/GLD, IEF early) — liquid ETFs, rebalanced with the book."),
 ("Constraints: ","long-only, no leverage, ~1% single-name cap, monthly rebalance."),
 ("Costs: ","15 bps round-trip applied to all turnover (~%d%%/yr); returns are net of these."%(V['ops']['avg_annual_turnover']*100)),
],size=12.5,gap=8)

# ===== 7. RISK FRAMEWORK =====
s=slide(); header(s,"Risk framework — the core of the strategy","RISK FRAMEWORK (15 PTS)")
kpis(s,1.5,[("VOLATILITY",pc(S['Vol']),f"vs market {pc(SPY['Vol'])}"),
            ("MAX DRAWDOWN",pc(S['MaxDD']),f"vs market {pc(SPY['MaxDD'])}"),
            ("MARKET BETA",f"{AL['betas']['Mkt-RF']:.2f}","low net exposure"),
            ("SORTINO",f"{S['Sortino']:.2f}" if S.get('Sortino') else "—","downside-adjusted")],w=2.75)
bullets(s,0.55,3.2,12.2,3.3,[
 ("Primary control — the macro overlay: ","de-grosses to as low as 30% equities in stress; this is what produces the −21% drawdown vs the market's −46/−50%."),
 ("Factor discipline: ","controlled exposures (market β %.2f, momentum +%.2f); alpha measured against FF5+Momentum with Newey-West standard errors."%(AL['betas']['Mkt-RF'],AL['betas']['Mom'])),
 ("Diversification & limits: ","~125 names, 11 sectors, ~1%% single-name cap, no leverage, no shorting — failure of any one name is immaterial."),
 ("Defensive convexity: ","the risk-off sleeve (Treasuries + gold) tends to rise when equities fall, cushioning crises further."),
],size=13.5,gap=9)

# ===== 8. BACK-TEST RESULTS =====
s=slide(); header(s,"Back-test results (simulated, 2002–2026)","SIMULATED · NET OF 15 BPS COSTS · GROSS OF FEES")
s.shapes.add_picture("fig_deck_equity.png",Inches(0.5),Inches(1.5),width=Inches(7.4))
txt(s,8.1,1.5,4.7,0.35,"PERFORMANCE SUMMARY",size=12,color=NAVY,bold=True)
table(s,8.1,1.95,4.7,0.52,["","CAGR","Shrp","MaxDD"],
 [["Fund (full)",pc(S['CAGR']),f"{S['Sharpe']:.2f}",pc(S['MaxDD'])],
  ["In-sample",pc(ISm['CAGR']),f"{ISm['Sharpe']:.2f}",pc(ISm['MaxDD'])],
  ["Out-of-sample",pc(OOS['CAGR']),f"{OOS['Sharpe']:.2f}",pc(OOS['MaxDD'])],
  ["S&P 500",pc(SPY['CAGR']),f"{SPY['Sharpe']:.2f}",pc(SPY['MaxDD'])]],colw=[1.5,1,0.9,1.1],fs=10.5)
txt(s,8.1,4.4,4.7,2.0,[[("Alpha +%.1f%% / yr (t=%.1f)"%(AL['annualized']*100,AL['tstat']),True,POS)],
  [("vs FF5+Momentum, Newey-West SEs.",False,GREY)],
  [(" ",False,GREY)],
  [("Out-of-sample Sharpe %.2f vs %.2f in-sample — a modest step-down, not the collapse overfitting produces; still well above the market OOS."%(OOS['Sharpe'],ISm['Sharpe']),False,NAVY)]],size=13)
txt(s,0.5,6.55,12,0.3,"Label: all figures SIMULATED on historical data; no live or paper track record is shown.",size=9,color=GREY)

# ===== 9. BACK-TEST HONESTY =====
s=slide(); header(s,"Back-test honesty — the red-flag checklist","BACK-TEST HONESTY (25 PTS)")
table(s,0.55,1.55,12.2,0.58,["Common red flag","How we handle it"],
 [["Survivorship bias","Measured on Bloomberg point-in-time S&P 500 (1,085 names incl. delisted): only ~%.1f%%/yr — quantified, not assumed."%(BZ['cagr_drag']*100)],
  ["Look-ahead","All signals lagged ≥1 month; trailing-only z-scores; 2000–01 reserved as warm-up."],
  ["No transaction costs","15 bps round-trip on all turnover; every figure is net of costs."],
  ["Overfitting","Train/test split (2002–15 vs 2016–26); OOS Sharpe close to IS (no collapse); standard enhancements tested and rejected."],
  ["Sharpe too good / undisclosed","Sharpe %.2f with full methodology; deflated-Sharpe and factor regression in the appendix."%(S['Sharpe'])],
  ["Simulated vs live mislabelled","Everything is labelled SIMULATED; the fund is fictional."]],colw=[1.5,3.5],fs=11.5)

# ===== 10. STRESS TESTS =====
s=slide(); header(s,"Stress tests — behaviour in every major crisis","RISK FRAMEWORK")
s.shapes.add_picture("fig_deck_stress.png",Inches(0.5),Inches(1.5),width=Inches(7.5))
bullets(s,8.2,1.7,4.6,4.6,[
 ("The pattern: ","in every major drawdown the overlay de-grosses and the defensive sleeve cushions, so we fall far less than the market."),
 ("GFC & COVID: ","the largest gaps — exactly when capital preservation matters most."),
 ("2022: ","bonds and stocks fell together; we disclose this regime risk, yet the out-of-sample period (incl. 2022) still beats the market."),
],size=13,gap=10)

# ===== 11. ROBUSTNESS =====
s=slide(); header(s,"Robustness — we tried to break it","RESEARCH PROCESS · ANTI-OVERFITTING")
s.shapes.add_picture("fig_bakeoff.png",Inches(0.45),Inches(1.5),width=Inches(7.5))
bullets(s,8.15,1.65,4.7,4.8,[
 ("Parameter stability: ","Sharpe %.2f–%.2f across 15 lookback × quantile settings — a plateau, not a spike."%(SEN['min_sharpe'],SEN['max_sharpe'])),
 ("Enhancements tested: ","risk-managed momentum, vol targeting, low-beta and quality tilts — none beat the simple design out-of-sample."),
 ("Discipline: ","we kept it simple by evidence, not by default. Complexity that overfits in-sample was rejected."),
],size=13,gap=11)

# ===== 12. CAPACITY =====
s=slide(); header(s,"Capacity — the $100M fits many times over","CAPACITY (QUANT BASIS)")
kpis(s,1.6,[("RAISE",f"${CAP['commitment_raise']/1e6:.0f}M","the commitment"),
            ("SOFT CAPACITY",f"${CAP['soft_cap_low']/1e9:.1f}–{CAP['soft_cap_high']/1e9:.1f}B","strategy-level"),
            ("HEADROOM",f"{CAP['headroom_low']:.0f}–{CAP['headroom_high']:.0f}×","vs the raise"),
            ("USED AT RAISE",f"~{CAP['pct_of_capacity_at_raise']*100:.0f}%","of capacity")],w=2.75)
bullets(s,0.55,3.3,12.2,3.0,[
 ("Basis: ","5%% of average daily volume over 2 days per name, across ~111 holdings. ADV %s."%CAP['adv_assumption']),
 ("Implication: ","the full $100M deploys immediately in daily-liquid names; we never approach our own market impact."),
 ("Defensive sleeve: ","held in the most liquid ETFs on earth (Treasuries, gold) — no capacity constraint."),
],size=14,gap=10)

# ===== 13. FUND TERMS =====
s=slide(); header(s,"Fund structure & terms — aligned with liquidity","FUND STRUCTURE & TERMS (10 PTS)")
txt(s,0.55,1.45,6.0,0.35,"TERMS",size=12,color=NAVY,bold=True)
table(s,0.55,1.9,6.0,0.55,["Term","What we charge"],
 [["Management fee","1.0% per annum"],["Performance fee","15% of return ABOVE the S&P 500"],
  ["High-water mark","Yes (relative to benchmark)"],["Liquidity","Monthly, no lock-up"],
  ["Crystallisation","Annual"]],colw=[1.3,2.2],fs=11)
txt(s,6.85,1.45,6.0,0.35,"WHAT THE INVESTOR RECEIVES (NET)",size=12,color=NAVY,bold=True)
table(s,6.85,1.9,6.0,0.55,["","CAGR","Sharpe","MaxDD"],
 [["Gross strategy",pc(NF['gross']['CAGR']),f"{NF['gross']['Sharpe']:.2f}",pc(NF['gross']['MaxDD'])],
  ["NET to investor",pc(NF['net']['CAGR']),f"{NF['net']['Sharpe']:.2f}",pc(NF['net']['MaxDD'])],
  ["S&P 500",pc(NF['spy']['CAGR']),f"{NF['spy']['Sharpe']:.2f}",pc(NF['spy']['MaxDD'])]],colw=[1.6,1,1,1],fs=11)
bullets(s,0.55,4.4,12.2,2.0,[
 ("Why not 2/20: ","the strategy is daily-liquid and partly factor-replicable, so we charge below hedge-fund terms."),
 ("Benchmark-relative performance fee: ","we are paid only on return ABOVE the S&P 500 — not on market beta you can buy for 5 bps."),
 ("Net of all fees: ","the investor still earns %s vs the market's %s, at higher Sharpe and ~half the drawdown."%(pc(NF['net']['CAGR']),pc(NF['spy']['CAGR']))),
],size=13.5,gap=9)

# ===== 14. WHY NOW =====
s=slide(); header(s,"Why now","WHY NOW")
bullets(s,0.55,1.7,12.2,4.6,[
 ("Regime uncertainty is elevated: ","rates, inflation and geopolitics keep drawdown risk high — a strategy built to de-risk systematically is timely."),
 ("Defensive assets pay again: ","with positive real yields, the Treasuries-and-gold sleeve earns a real return while waiting, unlike the zero-rate decade."),
 ("Factor crowding favours discipline over cleverness: ","raw signals are commoditised; durable edge now lives in risk management and execution, which is our core."),
 ("Liquidity and transparency are in demand: ","institutions want daily-liquid, rules-based, explainable strategies — exactly what we offer, with no black box."),
],size=15,gap=13)

# ===== 15. TEAM / AI / DISCLAIMER =====
s=slide(NAVY); rect(s,0,0,SW,0.06,GOLD)
txt(s,0.9,0.7,11.5,0.6,"Thank you — GSD2T Asset Management",size=30,color=WHITE,bold=True,font="Georgia")
txt(s,0.92,1.6,11.5,0.4,"Systematic. Diversified. Built to lose less.",size=15,color=GOLD)
txt(s,0.92,2.5,11.5,0.35,"TEAM",size=12,color=GOLD,bold=True)
txt(s,0.92,2.9,11.5,0.4,"[Names & roles — every member presents]",size=13,color=WHITE)
txt(s,0.92,3.6,11.5,0.35,"AI-TOOL DISCLOSURE",size=12,color=GOLD,bold=True)
txt(s,0.92,4.0,11.8,0.9,"AI assistants were used for code scaffolding, back-test engineering, document drafting and review. "
    "All strategy decisions, data choices and results were defined and verified by the team. No results were fabricated; "
    "all figures are reproducible from the committed code (see HOW_TO_RUN.md).",size=12,color=WHITE)
txt(s,0.92,5.4,11.5,0.35,"DATA & METHODS",size=12,color=GOLD,bold=True)
txt(s,0.92,5.8,11.8,0.6,"yfinance, Ken French Data Library, Bloomberg point-in-time S&P 500. Back-test 2002–2026, monthly, "
    "net of 15 bps costs. Full appendix in the Jupyter notebook.",size=12,color=WHITE)
txt(s,0.92,6.7,11.8,0.5,"Disclaimer: Fictional pitch for the ESADE Asset Management course. Not investment advice. All performance is "
    "SIMULATED; past performance does not indicate future results. No real fund names or live numbers are used.",size=9,color=GREY)

# ===== 16. APPENDIX — PORTFOLIO CONSTRUCTION ROBUSTNESS (Q&A backup) =====
s=slide(); header(s,"Appendix A — why we hold ~125 names, not 25","Q&A BACKUP · BEYOND THE 15-SLIDE PITCH")
txt(s,0.55,1.4,12.2,0.5,"\"Why not just hold 20–30 high-conviction names?\" — Because the concentrated book only wins in our single most-optimistic backtest, and reverses on the honest data.",size=13,color=NAVY,bold=True)
rows=[("Top 25 names","Top 25 names (high-conviction)","Top 25 names"),
      ("Top decile (~50)","Top decile (q=0.10, ~50)","Top decile (~50)"),
      ("Quartile (~125) — FLAGSHIP","Quartile (q=0.25) — V1 FLAGSHIP","Quartile (~125)"),
      ("Tercile (~165)","Tercile (q=0.33, ~165)","Tercile (~165)")]
trows=[[lab, f"{CV[ck]['full']['Sharpe']:.2f}", f"{CS['free'][sk]['Sharpe']:.2f}", f"{CS['free'][sk]['MaxDD']*100:.0f}%"] for lab,ck,sk in rows]
table(s,0.55,2.1,8.5,0.6,["Concentration","Monthly Sharpe\n(optimistic)","Survivorship-free\nSharpe (honest)","Surv-free\nMaxDD"],trows,colw=[2.4,1.3,1.4,1],fs=11)
txt(s,9.3,2.1,3.5,3.2,[[("THE PATTERN",True,NAVY)],
  [("Concentrated (25) wins only at monthly frequency on current constituents (Sharpe 1.20).",False,DARK)],
  [(" ",False,DARK)],
  [("On survivorship-free, point-in-time data it LOSES — quartile 0.54 vs 0.48, and −21% vs −31% drawdown.",False,DARK)]],size=12)
bullets(s,0.55,5.0,12.2,1.6,[
 ("Verdict: ","an edge that appears only in our most flattering configuration — and reverses on survivorship-free data and at quarterly frequency — is an artifact, not an edge. We pitch the diversified book: the number we can trust."),
 ("Also tested: ","inverse-volatility weighting lowered the Sharpe (0.99 vs 1.09); equal weight won on evidence."),
],size=12.5,gap=8)

# ===== 17. APPENDIX B — ASSUMPTION ROBUSTNESS (Q&A backup) =====
s=slide(); header(s,"Appendix B — every assumption stress-tested","Q&A BACKUP · BEYOND THE 15-SLIDE PITCH")
txt(s,0.55,1.4,12.2,0.5,"We measured or stress-tested every input. The two judgement calls — trading costs and the overlay calibration — both hold.",size=13,color=NAVY,bold=True)
# transaction-cost ladder
txt(s,0.55,2.05,5.7,0.35,"TRANSACTION COST — edge survives even at 6× the assumption",size=11,color=NAVY,bold=True)
tc_rows=[[f"{r['bps']} bps"+(" (current)" if r['bps']==15 else f" ({r['bps']//15}×)" if r['bps'] in (30,) else ""),
          pc(r['CAGR']),f"{r['Sharpe']:.2f}",pc(r['MaxDD'])] for r in RB['tcost']]
table(s,0.55,2.5,5.7,0.55,["Round-trip cost","CAGR","Sharpe","MaxDD"],tc_rows,colw=[1.7,1,1,1],fs=10.5)
txt(s,0.55,5.2,5.7,0.9,"Even at 100 bps — six times our assumption — Sharpe 0.82 and 11.3% CAGR still beat the market (0.60 / 10.0%). The 15 bps assumption is not load-bearing.",size=10.5,color=GREY)
# overlay grid
txt(s,6.7,2.05,6.1,0.35,"OVERLAY CALIBRATION — a plateau, not a curve-fit",size=11,color=NAVY,bold=True)
ov=RB['overlay']; head=["base \\ slope"]+[f"{sl:.3f}" for sl in ov['slopes']]
ov_rows=[[f"{b:.2f}"+(" ←" if abs(b-0.65)<1e-9 else "")]+[f"{ov['sharpe_grid'][i][j]:.2f}" for j in range(len(ov['slopes']))] for i,b in enumerate(ov['bases'])]
table(s,6.7,2.5,6.1,0.55,head,ov_rows,colw=[1.4,1,1,1],fs=10.5)
txt(s,6.7,5.0,6.1,1.1,f"Sharpe ranges only {ov['min']:.2f}–{ov['max']:.2f} across nine settings (baseline 55–75%, slope 0.125–0.225). "
    "The result doesn't hinge on the exact 0.65 / 0.175 — direct evidence we did not tune these to the backtest.",size=10.5,color=GREY)
bullets(s,0.55,6.1,12.2,0.7,[("Takeaway: ","every assumption in the model is either measured against data (ADV) or stress-tested (cost, overlay). None is a hidden knife-edge.")],size=12,gap=4)

prs.save("GSD2T_Pitch_Deck.pptx")
print(f"Saved GSD2T_Pitch_Deck.pptx ({len(prs.slides._sldIdLst)} slides)")
